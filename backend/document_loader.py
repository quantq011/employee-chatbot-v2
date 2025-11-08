"""
Document Loader Module
Handles extraction of text from various file formats: PDF, DOCX, PPTX, TXT, MD
"""
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import logging
 
# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
 
 
class DocumentLoader:
    """Load and extract text from various document formats"""
   
    SUPPORTED_FORMATS = ['.pdf', '.docx', '.pptx', '.txt', '.md']
   
    def __init__(self):
        """Initialize document loader with format handlers"""
        self.handlers = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx,
            '.txt': self._load_text,
            '.md': self._load_text,
        }
   
    def load_document(self, file_path: str) -> Dict[str, any]:
        """
        Load a document and extract its text content
       
        Args:
            file_path: Path to the document file
           
        Returns:
            Dictionary containing:
                - content: Extracted text
                - metadata: File information (name, type, size, etc.)
                - success: Boolean indicating success/failure
                - error: Error message if failed
        """
        path = Path(file_path)
       
        if not path.exists():
            return {
                'success': False,
                'error': f"File not found: {file_path}",
                'content': None,
                'metadata': None
            }
       
        extension = path.suffix.lower()
       
        if extension not in self.SUPPORTED_FORMATS:
            return {
                'success': False,
                'error': f"Unsupported format: {extension}. Supported: {self.SUPPORTED_FORMATS}",
                'content': None,
                'metadata': None
            }
       
        try:
            handler = self.handlers.get(extension)
            content = handler(path)
           
            metadata = {
                'filename': path.name,
                'file_type': extension,
                'file_size': path.stat().st_size,
                'file_path': str(path.absolute()),
            }
           
            logger.info(f"✓ Loaded {path.name} ({len(content)} characters)")
           
            return {
                'success': True,
                'content': content,
                'metadata': metadata,
                'error': None
            }
           
        except Exception as e:
            logger.error(f"✗ Error loading {path.name}: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': None,
                'metadata': None
            }
   
    def _load_pdf(self, path: Path) -> str:
        """Extract text from PDF using PyMuPDF"""
        try:
            import fitz  # PyMuPDF
           
            text_content = []
            doc = fitz.open(path)
           
            for page_num, page in enumerate(doc, 1):
                text = page.get_text()
                if text.strip():
                    text_content.append(text)
           
            doc.close()
           
            return '\n\n'.join(text_content)
           
        except ImportError:
            raise ImportError(
                "PyMuPDF (fitz) not installed. Install with: pip install PyMuPDF"
            )
   
    def _load_docx(self, path: Path) -> str:
        """Extract text from DOCX using python-docx"""
        try:
            from docx import Document
           
            doc = Document(path)
            text_content = []
           
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
           
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_content.append(row_text)
           
            return '\n\n'.join(text_content)
           
        except ImportError:
            raise ImportError(
                "python-docx not installed. Install with: pip install python-docx"
            )
   
    def _load_pptx(self, path: Path) -> str:
        """Extract text from PPTX using python-pptx"""
        try:
            from pptx import Presentation
           
            prs = Presentation(path)
            text_content = []
           
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
               
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
               
                if slide_text:
                    text_content.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_text))
           
            return '\n\n'.join(text_content)
           
        except ImportError:
            raise ImportError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
   
    def _load_text(self, path: Path) -> str:
        """Load plain text or markdown files"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise ValueError(f"Could not decode file with common encodings")
   
    def load_directory(self, directory_path: str, recursive: bool = True) -> List[Dict]:
        """
        Load all supported documents from a directory
       
        Args:
            directory_path: Path to directory
            recursive: Whether to search subdirectories
           
        Returns:
            List of document results (same format as load_document)
        """
        path = Path(directory_path)
       
        if not path.exists() or not path.is_dir():
            logger.error(f"Directory not found: {directory_path}")
            return []
       
        results = []
        pattern = '**/*' if recursive else '*'
       
        for file_path in path.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_FORMATS:
                result = self.load_document(str(file_path))
                results.append(result)
       
        logger.info(f"✓ Loaded {len(results)} documents from {directory_path}")
        return results
 
 
# Example usage
if __name__ == "__main__":
    loader = DocumentLoader()
   
    # Test loading a single document
    result = loader.load_document("example.pdf")
    if result['success']:
        print(f"Content length: {len(result['content'])}")
        print(f"Metadata: {result['metadata']}")
   
    # Test loading a directory
    results = loader.load_directory("./documents", recursive=True)
    print(f"Loaded {len(results)} documents")
 
 